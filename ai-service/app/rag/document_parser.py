from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from fastapi import status

from app.utils.errors import AppError


@dataclass(frozen=True)
class ParsedSection:
    text: str
    reference: str


class DocumentParser:
    supported_extensions = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".xls"}

    def parse(self, filename: str, content: bytes) -> list[ParsedSection]:
        extension = Path(filename).suffix.lower()
        if extension == ".pdf":
            return self._parse_pdf(content)
        if extension == ".docx":
            return self._parse_docx(content)
        if extension == ".pptx":
            return self._parse_pptx(content)
        if extension in {".xlsx", ".xlsm", ".xls"}:
            return self._parse_excel(content)
        raise AppError(
            "Unsupported file type. Upload PDF, DOCX, PPTX, or Excel files.",
            status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
            "UNSUPPORTED_KNOWLEDGE_FILE",
        )

    def _parse_pdf(self, content: bytes) -> list[ParsedSection]:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise AppError(
                "PDF parser dependency is not installed",
                status.HTTP_503_SERVICE_UNAVAILABLE,
                "PARSER_MISSING",
            ) from exc

        reader = PdfReader(BytesIO(content))
        sections = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                sections.append(ParsedSection(text=text, reference=f"page {index}"))
        return sections

    def _parse_docx(self, content: bytes) -> list[ParsedSection]:
        try:
            from docx import Document
        except ImportError as exc:
            raise AppError(
                "DOCX parser dependency is not installed",
                status.HTTP_503_SERVICE_UNAVAILABLE,
                "PARSER_MISSING",
            ) from exc

        document = Document(BytesIO(content))
        paragraphs = [
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        ]
        sections = []
        for index, paragraph in enumerate(paragraphs, start=1):
            sections.append(ParsedSection(text=paragraph, reference=f"paragraph {index}"))
        return sections

    def _parse_pptx(self, content: bytes) -> list[ParsedSection]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise AppError(
                "PPTX parser dependency is not installed",
                status.HTTP_503_SERVICE_UNAVAILABLE,
                "PARSER_MISSING",
            ) from exc

        presentation = Presentation(BytesIO(content))
        sections = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                sections.append(
                    ParsedSection(text="\n".join(texts), reference=f"slide {index}")
                )
        return sections

    def _parse_excel(self, content: bytes) -> list[ParsedSection]:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise AppError(
                "Excel parser dependency is not installed",
                status.HTTP_503_SERVICE_UNAVAILABLE,
                "PARSER_MISSING",
            ) from exc

        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        sections = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sections.append(
                    ParsedSection(text="\n".join(rows), reference=f"sheet {sheet.title}")
                )
        return sections


document_parser = DocumentParser()
